#%%
import pptx as pp

pp_file = '/Users/mrenke/Desktop/StressRisk/findings_01.pptx'
#%%
from pptx import Presentation
prs = Presentation(pp_file)




# %%
pati = '/Users/mrenke/data/ds-stressrisk/derivatives/cogmodels/figures/bambi-RNP-fit_model1/diffOFdiff_rnp.png'

# get reference to first shape in first slide
sp = prs.slides[0].shapes[0]

# add a picture shape to slide
pic = sld.shapes.add_picture(path, x, y, cx, cy)

